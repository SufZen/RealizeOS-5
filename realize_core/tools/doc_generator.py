"""
Document generation — create PDF, DOCX, PPTX files from content.

Uses:
- reportlab for PDF generation
- python-docx for Word documents
- python-pptx for PowerPoint presentations

Falls back gracefully if libraries aren't installed.
"""
import logging
from datetime import datetime
from pathlib import Path

logger = logging.getLogger(__name__)


def generate_pdf(
    title: str,
    content: str,
    output_dir: Path,
    author: str = "RealizeOS",
) -> dict:
    """
    Generate a PDF document from text content.

    Returns:
        {path, size, pages} or {error}
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import cm
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
    except ImportError:
        return {"error": "reportlab not installed. Run: pip install reportlab"}

    output_dir.mkdir(parents=True, exist_ok=True)
    safe_name = "".join(c if c.isalnum() or c in "-_ " else "" for c in title)[:50].strip()
    filename = f"{safe_name or 'document'}.pdf"
    filepath = output_dir / filename

    doc = SimpleDocTemplate(
        str(filepath),
        pagesize=A4,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
        leftMargin=2.5 * cm,
        rightMargin=2.5 * cm,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("DocTitle", parent=styles["Title"], fontSize=18, spaceAfter=20)
    body_style = ParagraphStyle("DocBody", parent=styles["Normal"], fontSize=11, leading=16)
    heading_style = ParagraphStyle("DocHeading", parent=styles["Heading2"], fontSize=14, spaceAfter=10, spaceBefore=15)

    elements = []
    elements.append(Paragraph(title, title_style))
    elements.append(Spacer(1, 12))

    # Parse content into paragraphs (handle markdown headers)
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            elements.append(Spacer(1, 8))
        elif line.startswith("## "):
            elements.append(Paragraph(line[3:], heading_style))
        elif line.startswith("# "):
            elements.append(Paragraph(line[2:], title_style))
        elif line.startswith("- "):
            elements.append(Paragraph(f"\u2022 {line[2:]}", body_style))
        else:
            # Escape XML special chars for reportlab
            safe_line = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            elements.append(Paragraph(safe_line, body_style))

    doc.build(elements)

    return {
        "path": str(filepath),
        "filename": filename,
        "size": filepath.stat().st_size,
        "format": "pdf",
    }


def generate_docx(
    title: str,
    content: str,
    output_dir: Path,
    author: str = "RealizeOS",
) -> dict:
    """
    Generate a Word document from text content.

    Returns:
        {path, size} or {error}
    """
    try:
        from docx import Document
        from docx.shared import Pt  # noqa: F401
    except ImportError:
        return {"error": "python-docx not installed. Run: pip install python-docx"}

    output_dir.mkdir(parents=True, exist_ok=True)
    safe_name = "".join(c if c.isalnum() or c in "-_ " else "" for c in title)[:50].strip()
    filename = f"{safe_name or 'document'}.docx"
    filepath = output_dir / filename

    doc = Document()
    doc.core_properties.author = author
    doc.core_properties.title = title

    doc.add_heading(title, level=0)

    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:], style="List Bullet")
        else:
            doc.add_paragraph(line)

    doc.save(str(filepath))

    return {
        "path": str(filepath),
        "filename": filename,
        "size": filepath.stat().st_size,
        "format": "docx",
    }


def generate_pptx(
    title: str,
    content: str,
    output_dir: Path,
    author: str = "RealizeOS",
) -> dict:
    """
    Generate a PowerPoint presentation from text content.

    Splits content by ## headings into slides.

    Returns:
        {path, size, slides} or {error}
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt  # noqa: F401
    except ImportError:
        return {"error": "python-pptx not installed. Run: pip install python-pptx"}

    output_dir.mkdir(parents=True, exist_ok=True)
    safe_name = "".join(c if c.isalnum() or c in "-_ " else "" for c in title)[:50].strip()
    filename = f"{safe_name or 'presentation'}.pptx"
    filepath = output_dir / filename

    prs = Presentation()
    prs.core_properties.author = author
    prs.core_properties.title = title

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = f"Generated by RealizeOS — {datetime.now().strftime('%Y-%m-%d')}"

    # Split content into slides by ## headings
    current_heading = ""
    current_body = []
    slide_count = 1

    for line in content.split("\n"):
        line = line.strip()
        if line.startswith("## "):
            if current_heading or current_body:
                _add_content_slide(prs, current_heading, current_body)
                slide_count += 1
            current_heading = line[3:]
            current_body = []
        elif line.startswith("# "):
            continue  # Skip top-level headers (already in title slide)
        elif line:
            current_body.append(line)

    # Add last section
    if current_heading or current_body:
        _add_content_slide(prs, current_heading, current_body)
        slide_count += 1

    prs.save(str(filepath))

    return {
        "path": str(filepath),
        "filename": filename,
        "size": filepath.stat().st_size,
        "slides": slide_count,
        "format": "pptx",
    }


def _add_content_slide(prs, heading: str, body_lines: list[str]):
    """Add a content slide to the presentation."""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    if heading:
        slide.shapes.title.text = heading
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines[:8]):  # Max 8 bullets per slide
        if i == 0:
            tf.text = line.lstrip("- ")
        else:
            p = tf.add_paragraph()
            p.text = line.lstrip("- ")


def get_available_formats() -> list[dict]:
    """Check which document formats are available."""
    formats = []

    try:
        import reportlab  # noqa: F401
        formats.append({"format": "pdf", "available": True, "library": "reportlab"})
    except ImportError:
        formats.append({"format": "pdf", "available": False, "library": "reportlab", "install": "pip install reportlab"})

    try:
        import docx  # noqa: F401
        formats.append({"format": "docx", "available": True, "library": "python-docx"})
    except ImportError:
        formats.append({"format": "docx", "available": False, "library": "python-docx", "install": "pip install python-docx"})

    try:
        import pptx  # noqa: F401
        formats.append({"format": "pptx", "available": True, "library": "python-pptx"})
    except ImportError:
        formats.append({"format": "pptx", "available": False, "library": "python-pptx", "install": "pip install python-pptx"})

    return formats
